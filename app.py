# -*- coding: utf-8 -*-
"""
Created on Thu Mar 16 10:35:49 2023

@author: Noura Aye
"""

from flask import Flask, request, render_template, send_file
import os
import io
import openpyxl
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/generate_slides', methods=['POST'])
def generate_slides():
    # Get the uploaded file
    excel_file = request.files['excel_file']
    
    # Parse the Excel file
    wb = openpyxl.load_workbook(excel_file)
    sheet = wb.active
    data = []
    for row in sheet.iter_rows(min_row=2, values_only=True):
        data.append(row)
    
    # Generate the PowerPoint slides
    prs = Presentation('template.pptx')
    for row in data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = row[0]
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = row[1]
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        #subtitle = slide.placeholders[1]
        #subtitle.text = row[1]
        #subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        #bullet_slide = slide.placeholders[2]
        #bullet_slide.text = row[2]
        #bullet_slide.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        #bullet_slide.text_frame.paragraphs[0].font.size = Cm(0.5)
    
    # Save the PowerPoint file to a BytesIO object
    pptx_data = io.BytesIO()
    prs.save(pptx_data)
    pptx_data.seek(0)
    
    # Send the PowerPoint file to the user for download
    # response = send_file(
    #     pptx_data,
    #     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
    #     as_attachment=True
    # )
    # response.headers["Content-Disposition"] = "attachment; filename=generated_slides.pptx"
    # return response
    return send_file(
        pptx_data,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        as_attachment=True,
        download_name='generated_slides.pptx'
    )

if __name__ == '__main__':
    app.run(debug=True)
